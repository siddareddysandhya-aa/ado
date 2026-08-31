from io import BytesIO

import pandas as pd
import plotly.express as px
from pptx import Presentation

from ado_agile_metrics.reports import powerpoint_bytes
from ado_agile_metrics.storage import SnapshotStore


def test_powerpoint_includes_one_slide_per_chart():
    figure = px.line(x=[pd.Timestamp("2026-08-12"), pd.Timestamp("2026-08-13")], y=[3, 5], title="Velocity Trend")

    presentation = Presentation(BytesIO(powerpoint_bytes("Summary", {"velocity": 5.0}, ["Delivery improved."], [figure])))

    assert len(presentation.slides) == 4
    assert any(shape.has_text_frame and shape.text == "Engineering Metrics" for shape in presentation.slides[1].shapes)


def test_powerpoint_snapshot_uses_requested_headline_metrics():
    presentation = Presentation(BytesIO(powerpoint_bytes("Summary", {"velocity": 5.0, "predictability": 99.8, "lead_time": 3.0, "cycle_time": 2.0}, [], [])))
    snapshot_text = " ".join(shape.text for shape in presentation.slides[2].shapes if shape.has_text_frame)

    assert "PREDICTABILITY" not in snapshot_text
    assert "VELOCITY" in snapshot_text
    assert "LEAD TIME" in snapshot_text
    assert "CYCLE TIME" in snapshot_text


def test_powerpoint_includes_selected_area_path_context():
    presentation = Presentation(BytesIO(powerpoint_bytes("Summary", {}, [], context={"area_path": "Day_Of_Operations\\ECS", "squad": "Jetstream"})))
    cover_text = " ".join(shape.text for shape in presentation.slides[0].shapes if shape.has_text_frame)

    assert "Area Path: Day_Of_Operations\\ECS" in cover_text
    assert "Squad: Jetstream" in cover_text


def test_saved_dashboard_can_be_reloaded(tmp_path):
    store = SnapshotStore(tmp_path / "metrics.db")
    configuration = {"iteration_filter": ["Sprint 1"], "metric_filter": ["Velocity"]}

    store.save_dashboard("Sprint review", configuration)

    assert store.dashboard_names() == ["Sprint review"]
    assert store.load_dashboard("Sprint review") == configuration


def test_last_used_dashboard_is_not_shown_as_a_named_preset(tmp_path):
    store = SnapshotStore(tmp_path / "metrics.db")

    store.save_last_used({"iteration_filter": ["Sprint 2"]})

    assert store.load_last_used() == {"iteration_filter": ["Sprint 2"]}
    assert store.dashboard_names() == []