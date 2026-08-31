"""Leadership report exports and SMTP delivery."""

from email.message import EmailMessage
from io import BytesIO
import json
import os
import smtplib
from typing import Iterable

import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen.canvas import Canvas


def pdf_bytes(title: str, metrics: dict[str, float], insights: list[str]) -> bytes:
    """Generate a concise leadership-ready PDF summary."""
    output = BytesIO()
    canvas = Canvas(output, pagesize=letter)
    canvas.setTitle(title)
    canvas.setFont("Helvetica-Bold", 18)
    canvas.drawString(54, 750, title)
    canvas.setFont("Helvetica", 11)
    y_position = 715
    for name, value in metrics.items():
        canvas.drawString(54, y_position, f"{name.replace('_', ' ').title()}: {value:.1f}")
        y_position -= 20
    y_position -= 12
    canvas.setFont("Helvetica-Bold", 12)
    canvas.drawString(54, y_position, "Key observations")
    canvas.setFont("Helvetica", 11)
    for insight in insights:
        y_position -= 20
        canvas.drawString(54, y_position, f"- {insight}"[:115])
    canvas.save()
    return output.getvalue()


NAVY = RGBColor(18, 35, 56)
BLUE = RGBColor(0, 112, 192)
TEAL = RGBColor(0, 148, 136)
CORAL = RGBColor(230, 103, 89)
INK = RGBColor(36, 48, 65)
MUTED = RGBColor(99, 115, 129)
PALE_BLUE = RGBColor(232, 242, 250)
WHITE = RGBColor(255, 255, 255)
CHART_COLORS = ("#1677C8", "#16826C", "#D9604A", "#7768A8", "#C99222")


def _text(slide, value: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = value
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _header(slide, title: str, page: int) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    _text(slide, title, 0.62, 0.35, 11, 0.4, 24, NAVY, bold=True)
    _text(slide, "AGILE METRICS PORTAL", 0.62, 7.08, 4, 0.2, 8, MUTED, bold=True)
    _text(slide, str(page), 12.0, 7.08, 0.6, 0.2, 8, MUTED, align=PP_ALIGN.RIGHT)


def _metric_card(slide, name: str, value: float, left: float, accent: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.48), Inches(3.76), Inches(1.85))
    card.fill.solid()
    card.fill.fore_color.rgb = PALE_BLUE
    card.line.color.rgb = PALE_BLUE
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(4.48), Inches(0.12), Inches(1.85))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()
    _text(slide, name.replace("_", " ").upper(), left + 0.33, 4.81, 3.1, 0.24, 10, MUTED, bold=True)
    _text(slide, f"{value:.1f}", left + 0.33, 5.22, 3.1, 0.5, 26, INK, bold=True)


def _style_chart(figure: go.Figure) -> None:
    """Apply a readable executive palette before a Plotly figure becomes a slide image."""
    figure.update_layout(
        title=None,
        template="plotly_white",
        paper_bgcolor="white",
        plot_bgcolor="white",
        colorway=CHART_COLORS,
        font={"family": "Aptos", "color": "#243041"},
        margin={"l": 65, "r": 35, "t": 35, "b": 55},
    )
    for index, trace in enumerate(figure.data):
        color = CHART_COLORS[index % len(CHART_COLORS)]
        if trace.type == "pie":
            trace.update(marker={"colors": CHART_COLORS})
        elif trace.type == "scatter":
            trace.update(marker={"color": color}, line={"color": color})
        else:
            trace.update(marker={"color": color})


def _chart_image(figure: go.Figure, width: int, height: int) -> BytesIO:
    """Render a styled chart image that is safe to embed in a presentation."""
    chart_figure = go.Figure(json.loads(pio.to_json(figure)))
    _style_chart(chart_figure)
    return BytesIO(chart_figure.to_image(format="png", width=width, height=height, scale=1))


def _figure_title(figure: go.Figure) -> str:
    """Return a stable display title for chart selection and slide labeling."""
    return str(figure.layout.title.text or "Agile metric")


def _add_chart_panel(slide, figure: go.Figure, left: float, top: float, width: float, height: float) -> None:
    """Place a chart inside an aligned, lightly framed reporting panel."""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALE_BLUE
    panel.line.color.rgb = RGBColor(214, 227, 239)
    _text(slide, _figure_title(figure), left + 0.2, top + 0.16, width - 0.4, 0.23, 11, NAVY, bold=True)
    slide.shapes.add_picture(_chart_image(figure, 900, 480), Inches(left + 0.12), Inches(top + 0.42), width=Inches(width - 0.24), height=Inches(height - 0.55))


def powerpoint_bytes(title: str, metrics: dict[str, float], insights: list[str], figures: Iterable[go.Figure] = (), context: dict[str, str] | None = None) -> bytes:
    """Generate a reference-style sprint metrics deck plus chart appendix slides."""
    context = context or {}
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    cover = presentation.slides.add_slide(blank)
    fill = cover.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    accent = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.35), Inches(0.14), Inches(3.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    _text(cover, "DAY OF OPERATIONS", 1.12, 1.45, 6, 0.3, 12, RGBColor(126, 211, 201), bold=True)
    _text(cover, "Engineering Metrics", 1.12, 1.96, 10.6, 0.7, 34, WHITE, bold=True)
    _text(cover, title, 1.12, 2.85, 9.5, 0.35, 16, RGBColor(202, 216, 232))
    _text(cover, "Sprint delivery review", 1.12, 3.35, 7, 0.35, 16, RGBColor(202, 216, 232))
    _text(cover, f"Project: {context.get('project', 'All selected projects')}", 1.12, 4.35, 8.8, 0.28, 12, WHITE)
    _text(cover, f"Squad: {context.get('squad', 'All selected squads')}", 1.12, 4.75, 8.8, 0.28, 12, WHITE)
    _text(cover, f"Area Path: {context.get('area_path', 'All selected areas')}", 1.12, 5.15, 10.5, 0.28, 12, WHITE)
    _text(cover, f"Sprint: {context.get('sprint', 'Selected iterations')}", 1.12, 5.55, 8.8, 0.28, 12, WHITE)
    _text(cover, "Azure DevOps Agile Metrics", 1.12, 6.62, 7, 0.3, 11, RGBColor(202, 216, 232))

    chart_figures = list(figures)
    metric_figures = [figure for figure in chart_figures if any(label in _figure_title(figure) for label in ("Velocity", "Burndown", "Lead Time", "Cycle Time"))][:4]
    if not metric_figures:
        metric_figures = chart_figures[:4]
    remaining_figures = [figure for figure in chart_figures if figure not in metric_figures]
    scorecard = presentation.slides.add_slide(blank)
    _header(scorecard, "Engineering Metrics", 2)
    positions = ((0.55, 1.0, 6.05, 2.75), (6.73, 1.0, 6.05, 2.75), (0.55, 3.88, 6.05, 2.75), (6.73, 3.88, 6.05, 2.75))
    for figure, position in zip(metric_figures, positions):
        _add_chart_panel(scorecard, figure, *position)

    summary = presentation.slides.add_slide(blank)
    _header(summary, "Sprint Summary", 3)
    _text(summary, f"{context.get('squad', 'All selected squads')} | {context.get('area_path', 'All selected areas')} | {context.get('sprint', 'Selected iterations')}", 0.65, 0.8, 11.7, 0.22, 10, MUTED)
    headline_metrics = [(name, metrics.get(name, 0.0)) for name in ("velocity", "lead_time", "cycle_time")]
    for index, (name, value) in enumerate(headline_metrics):
        _metric_card(summary, name, value, 0.6 + index * 4.18, (BLUE, TEAL, CORAL)[index])
    _text(summary, "Delivery highlights", 0.65, 1.25, 5.5, 0.3, 16, NAVY, bold=True)
    for index, insight in enumerate(insights[:5]):
        _text(summary, f"{index + 1}. {insight}", 0.75, 1.85 + index * 0.48, 11.7, 0.35, 13, INK)

    closing = presentation.slides.add_slide(blank)
    fill = closing.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    _text(closing, "Thank you", 0.8, 2.85, 11.7, 0.7, 36, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text(closing, "Engineering Metrics", 0.8, 3.7, 11.7, 0.35, 16, RGBColor(202, 216, 232), align=PP_ALIGN.CENTER)

    for page, figure in enumerate(remaining_figures, start=5):
        chart_slide = presentation.slides.add_slide(blank)
        _header(chart_slide, _figure_title(figure), page)
        _add_chart_panel(chart_slide, figure, 0.55, 1.05, 12.23, 5.75)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def send_email_report(recipient: str, subject: str, insights: list[str], attachment: bytes) -> None:
    """Send a PDF leadership report through explicitly configured SMTP."""
    host = os.getenv("SMTP_HOST", "")
    sender = os.getenv("SMTP_FROM", "")
    if not host or not sender:
        raise ValueError("Set SMTP_HOST and SMTP_FROM before sending email reports.")
    message = EmailMessage()
    message["From"] = sender
    message["To"] = recipient
    message["Subject"] = subject
    message.set_content("\n".join(insights))
    message.add_attachment(attachment, maintype="application", subtype="pdf", filename="agile-metrics-summary.pdf")
    with smtplib.SMTP(host, int(os.getenv("SMTP_PORT", "587")), timeout=30) as smtp:
        if os.getenv("SMTP_STARTTLS", "true").lower() == "true":
            smtp.starttls()
        username = os.getenv("SMTP_USERNAME")
        password = os.getenv("SMTP_PASSWORD")
        if username and password:
            smtp.login(username, password)
        smtp.send_message(message)